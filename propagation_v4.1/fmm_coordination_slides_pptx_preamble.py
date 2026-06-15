"""
fmm_coordination_slides_pptx_preamble.py
========================================
FMM division coordination slides — PowerPoint chrome and helpers.

Pairs with: fmm_coordination_slides_pptx_skill.md (read that skill first).
This is the PowerPoint counterpart of fmm_coordination_slides_preamble.tex
(Beamer) and fmm_word_preamble.js (Word). It carries the *chrome* only —
palette, fonts, full-bleed backgrounds, and thin slide scaffolds. It does NOT
dictate content or lock layouts: build custom slides directly with python-pptx
when a slide needs something the scaffolds don't cover.

Colors mirror andreas_palette.md. Never edit the COLORS values independently;
update the palette and sweep all consuming files in one atomic commit.

Requires: python-pptx  (pip install python-pptx)

Last updated: May 28, 2026 — v2.0.0 (v3.2: new file. PowerPoint chrome/helpers
paired with the new fmm_coordination_slides_pptx_skill.md, formalizing the
sample validated with Andrea. Carries the v2.0.0 floor under the two-tracker
convention; kit-wide marker reset and kit package version land at v3.2 cycle
close.)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette (mirrors andreas_palette.md; hex without '#') ---------------------
COLORS = {
    "PrimaryBlue": "196E8C", "SecondaryBlue": "5E90A8", "LightBlue": "CBD8DF",
    "DarkBlue": "004B70", "SoftBlue": "ADD8E6", "DarkGreen": "4C7234",
    "LightGreen": "B4C896", "WarmCream": "FFF3D6",
    "DarkYellow": "F7E272", "MintGreen": "9FD0A3", "BrightBlue": "009ADE",
    "BrightYellow": "FFDA00", "BrightGreen": "8DBA38", "DeepGreen": "308144",
    "LabelGray": "A6A6A8", "White": "FFFFFF", "Black": "000000",
}
FONT = "Arial"  # FMM institutional font (matches the Word documents)

# --- Assets (must be in ASSET_DIR alongside the generated .pptx) ---------------
BG_TITLE = "BackgroundTitle.png"        # dark-blue title background (logos baked in)
BG_BODY_BLUE = "BackgroundBody.png"     # light-blue body background (logos baked in)
BG_BODY_WHITE = "BackgroundBodyWhite.png"  # white body background (logos baked in)

# --- Geometry / safe zones -----------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.9)
# Logos are baked into the bottom corners of every background; keep content above.
LOGO_SAFE_BOTTOM = Inches(6.7)
# On BackgroundTitle the right ~third brightens to cyan; DarkYellow title text
# must stay in the dark left zone or it washes out.
TITLE_SAFE_WIDTH = Inches(8.5)


def color(name_or_hex):
    """RGBColor from a palette name or a raw 6-char hex string."""
    return RGBColor.from_string(COLORS.get(name_or_hex, name_or_hex))


def new_deck():
    """A blank 16:9 deck."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def background(slide, image_name, asset_dir="."):
    """Lay a full-bleed background image (logos are baked into the asset)."""
    slide.shapes.add_picture(os.path.join(asset_dir, image_name), 0, 0,
                             width=SLIDE_W, height=SLIDE_H)


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    """A zero-margin textbox (so alignment is predictable)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def run(paragraph, text, size, color_name, bold=False, font=FONT):
    """Add a styled run to a paragraph."""
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color(color_name)
    return r


# --- Slide scaffolds -----------------------------------------------------------
# Each returns the slide so the caller can keep adding content. They set the
# canvas and the load-bearing chrome (background, title placement/colour); the
# caller supplies the body. Drop to raw python-pptx for anything custom.

def title_slide(prs, title, subtitle="", asset_dir="."):
    """Dark title slide: DarkYellow title + white subtitle, placed in the dark
    left zone (where DarkYellow reads; it would wash out toward the bright edge)."""
    s = _blank(prs)
    background(s, BG_TITLE, asset_dir)
    tf = textbox(s, MARGIN, Inches(2.3), TITLE_SAFE_WIDTH, Inches(2.4))
    run(tf.paragraphs[0], title, 44, "DarkYellow", bold=True)
    if subtitle:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        run(p, subtitle, 22, "White")
    return s


def body_slide(prs, title, background_choice="white", asset_dir="."):
    """Body slide with a PrimaryBlue frame title. background_choice: 'white' or
    'blue'. Returns (slide, content_top_inches) — add body content below
    content_top and keep it above LOGO_SAFE_BOTTOM."""
    s = _blank(prs)
    bg = BG_BODY_WHITE if background_choice == "white" else BG_BODY_BLUE
    background(s, bg, asset_dir)
    tf = textbox(s, MARGIN, Inches(0.6), SLIDE_W - 2 * MARGIN, Inches(1.0))
    run(tf.paragraphs[0], title, 34, "PrimaryBlue", bold=True)
    return s, 1.9  # caller starts content around 1.9in


def bullets(slide, items, left=MARGIN, top=Inches(1.9), width=Inches(11.0),
            height=Inches(4.0), marker_color="PrimaryBlue", size=18):
    """Bold-lead-in bullets. items: list of (lead, rest) tuples or plain str."""
    tf = textbox(slide, left, top, width, height)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(14)
        run(p, "\u25AA  ", size, marker_color, bold=True)
        if isinstance(item, tuple):
            lead, rest = item
            run(p, lead, size, "DarkBlue", bold=True)
            run(p, rest, size, "Black")
        else:
            run(p, item, size, "Black")
    return tf


def callout(slide, text, left, top, width, height=Inches(0.8), fill="BrightBlue"):
    """A small accent callout box (BrightBlue by default — use sparingly)."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color(fill)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    run(tf.paragraphs[0], text, 16, "White", bold=True)
    return box


def action_slide(prs, title, rows, background_choice="white", asset_dir="."):
    """Action-item slide: status dot + task + owner. rows: list of
    (status_color_name, task, owner). Defaults to the white background for
    legible owner labels; on the blue background LabelGray owners read faint
    (acceptable when owners are deliberately de-emphasized)."""
    s, _ = body_slide(prs, title, background_choice, asset_dir)
    y = 2.0
    for dot, task, owner in rows:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y + 0.07),
                               Inches(0.28), Inches(0.28))
        c.fill.solid(); c.fill.fore_color.rgb = color(dot); c.line.fill.background()
        run(textbox(s, Inches(1.5), Inches(y), Inches(7.5), Inches(0.6)).paragraphs[0],
            task, 20, "DarkBlue", bold=True)
        run(textbox(s, Inches(9.2), Inches(y + 0.03), Inches(3.0), Inches(0.6)).paragraphs[0],
            owner, 14, "LabelGray")
        y += 0.95
    return s


if __name__ == "__main__":
    # Self-test / usage example: rebuilds the validated 3-slide sample.
    prs = new_deck()
    title_slide(prs, "FMM Coordination Meeting", "Sample deck — palette & background check")
    s, _ = body_slide(prs, "Key points", "white")
    bullets(s, [
        ("Quarterly delivery on track. ", "All three knowledge products shipped before the division deadline."),
        ("Pipeline rebalanced. ", "Two notes moved to next quarter to free review capacity."),
        ("One item needs a decision. ", "Co-authorship on the regional brief is still open."),
    ])
    callout(s, "Decision needed by Friday", Inches(0.9), Inches(5.5), Inches(4.6))
    action_slide(prs, "Action items", [
        ("DeepGreen", "Finalize Q2 synthesis note", "Marta — done"),
        ("BrightBlue", "Circulate regional brief draft", "Andrea — this week"),
        ("LabelGray", "Schedule next coordination call", "Ops — pending date"),
    ], background_choice="blue")
    prs.save("module_selftest.pptx")
    print("module_selftest.pptx written")
